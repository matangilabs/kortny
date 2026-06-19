"""Tests for Office-format (DOCX/XLSX/PPTX) text extraction in slack_file_read."""

from __future__ import annotations

import io
import struct
import zipfile
from pathlib import Path

import pytest

import kortny.tools.slack_file_read as sfr
from kortny.tools.slack_file_read import (
    TextExtraction,
    _check_office_zip_bomb,
    _extract_docx_text,
    _extract_pptx_text,
    _extract_text,
    _extract_xlsx_text,
)

# ---------------------------------------------------------------------------
# Helpers to build minimal Office files in-memory
# ---------------------------------------------------------------------------


def make_docx(tmp_path: Path, paragraph: str, table_rows: list[list[str]]) -> Path:
    """Build a .docx with one paragraph and one table."""
    import docx

    doc = docx.Document()
    doc.add_paragraph(paragraph)
    if table_rows:
        tbl = doc.add_table(rows=len(table_rows), cols=len(table_rows[0]))
        for r_idx, row_data in enumerate(table_rows):
            for c_idx, cell_text in enumerate(row_data):
                tbl.rows[r_idx].cells[c_idx].text = cell_text
    path = tmp_path / "test.docx"
    doc.save(str(path))
    return path


def make_xlsx(tmp_path: Path, sheet_name: str, rows: list[list[object]]) -> Path:
    """Build a .xlsx with one sheet and given rows."""
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet_name
    for row in rows:
        ws.append(row)
    path = tmp_path / "test.xlsx"
    wb.save(str(path))
    return path


def make_pptx(tmp_path: Path, title: str, body: str) -> Path:
    """Build a .pptx with one slide containing a title text box and a body text box."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    title_box.text_frame.text = title

    body_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(3))
    body_box.text_frame.text = body

    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


def make_zip_bomb_bytes(*, entry_count: int = 1, file_size: int) -> bytes:
    """Build a ZIP whose infolist reports large file_size without storing real data.

    We craft a central directory entry that claims file_size bytes but the
    actual compressed data is a tiny stored entry. This exercises the
    metadata-only check in _check_office_zip_bomb.
    """
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_STORED) as zf:
        for i in range(entry_count):
            zf.writestr(f"entry_{i}.txt", b"x")
    raw = bytearray(buf.getvalue())

    # Patch the uncompressed size field in the central directory.
    # Find the central directory signature (PK\x01\x02).
    cd_sig = b"PK\x01\x02"
    pos = raw.find(cd_sig)
    if pos != -1:
        # Uncompressed size is at offset +24 from the CD signature (4 bytes, LE).
        struct.pack_into("<I", raw, pos + 24, file_size)

    return bytes(raw)


# ---------------------------------------------------------------------------
# DOCX tests
# ---------------------------------------------------------------------------


def test_docx_extracts_paragraph_and_table(tmp_path: Path) -> None:
    path = make_docx(
        tmp_path,
        paragraph="Hello from DOCX",
        table_rows=[["Alpha", "Beta"], ["Gamma", "Delta"]],
    )
    content = path.read_bytes()
    result = _extract_text(
        path,
        content=content,
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        max_chars=10_000,
    )
    assert result.supported is True
    assert result.backend == "docx"
    assert "Hello from DOCX" in (result.text or "")
    assert "Alpha" in (result.text or "")
    assert "Beta" in (result.text or "")
    assert "Gamma" in (result.text or "")


def test_docx_detected_by_extension(tmp_path: Path) -> None:
    path = make_docx(tmp_path, paragraph="Extension detect", table_rows=[])
    content = path.read_bytes()
    # Use a generic mime type but .docx extension — should still extract
    result = _extract_text(
        path, content=content, mime_type="application/octet-stream", max_chars=10_000
    )
    assert result.supported is True
    assert result.backend == "docx"


# ---------------------------------------------------------------------------
# XLSX tests
# ---------------------------------------------------------------------------


def test_xlsx_extracts_sheet_name_and_cells(tmp_path: Path) -> None:
    path = make_xlsx(
        tmp_path,
        sheet_name="Financials",
        rows=[["Revenue", "Cost"], [1000, 800], [2000, 1500]],
    )
    content = path.read_bytes()
    result = _extract_text(
        path,
        content=content,
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        max_chars=10_000,
    )
    assert result.supported is True
    assert result.backend == "xlsx"
    text = result.text or ""
    assert "Financials" in text
    assert "Revenue" in text
    assert "1000" in text


def test_xlsx_detected_by_extension(tmp_path: Path) -> None:
    path = make_xlsx(tmp_path, sheet_name="Sheet1", rows=[["A", "B"]])
    content = path.read_bytes()
    result = _extract_text(
        path, content=content, mime_type="application/octet-stream", max_chars=10_000
    )
    assert result.supported is True
    assert result.backend == "xlsx"


# ---------------------------------------------------------------------------
# PPTX tests
# ---------------------------------------------------------------------------


def test_pptx_extracts_slide_content(tmp_path: Path) -> None:
    path = make_pptx(tmp_path, title="Q1 Results", body="Revenue up 20%")
    content = path.read_bytes()
    result = _extract_text(
        path,
        content=content,
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        max_chars=10_000,
    )
    assert result.supported is True
    assert result.backend == "pptx"
    text = result.text or ""
    assert "Slide 1" in text
    assert "Q1 Results" in text
    assert "Revenue up 20%" in text


def test_pptx_detected_by_extension(tmp_path: Path) -> None:
    path = make_pptx(tmp_path, title="Title", body="Body")
    content = path.read_bytes()
    result = _extract_text(
        path, content=content, mime_type="application/octet-stream", max_chars=10_000
    )
    assert result.supported is True
    assert result.backend == "pptx"


# ---------------------------------------------------------------------------
# Zip-bomb guard tests
# ---------------------------------------------------------------------------


def test_zip_bomb_guard_triggers_on_large_uncompressed_size(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """A ZIP whose CD claims file_size > cap triggers the guard."""
    monkeypatch.setattr(sfr, "MAX_OFFICE_UNCOMPRESSED_BYTES", 100)
    # Build a zip that declares 101 bytes uncompressed
    bomb_bytes = make_zip_bomb_bytes(file_size=101)
    result = _check_office_zip_bomb(bomb_bytes)
    assert result is not None
    assert result.supported is False
    assert result.recoverable is True
    assert "office_file_exceeds_safe_limits" in result.warnings


def test_zip_bomb_guard_passes_normal_file(tmp_path: Path) -> None:
    """A normal DOCX passes the zip-bomb check (returns None)."""
    path = make_docx(tmp_path, paragraph="Normal doc", table_rows=[])
    content = path.read_bytes()
    result = _check_office_zip_bomb(content)
    assert result is None


def test_zip_bomb_guard_rejects_invalid_zip() -> None:
    """Non-ZIP bytes return unsupported (not recoverable)."""
    result = _check_office_zip_bomb(b"not a zip file at all")
    assert result is not None
    assert result.supported is False
    assert result.recoverable is False


def test_zip_bomb_guard_triggers_on_too_many_entries(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """A ZIP with more entries than MAX_OFFICE_ZIP_ENTRIES triggers the guard."""
    monkeypatch.setattr(sfr, "MAX_OFFICE_ZIP_ENTRIES", 2)
    bomb_bytes = make_zip_bomb_bytes(entry_count=3, file_size=10)
    result = _check_office_zip_bomb(bomb_bytes)
    assert result is not None
    assert result.supported is False
    assert result.recoverable is True
    assert "office_file_exceeds_safe_limits" in result.warnings


# ---------------------------------------------------------------------------
# Regression: existing types still carry backend label and same text
# ---------------------------------------------------------------------------


def test_csv_extraction_has_backend_label(tmp_path: Path) -> None:
    """CSV extraction is byte-identical to before and now carries backend='text'."""
    csv_content = b"ticker,price\nAAPL,200\n"
    path = tmp_path / "data.csv"
    path.write_bytes(csv_content)
    result = _extract_text(
        path, content=csv_content, mime_type="text/csv", max_chars=10_000
    )
    assert result.supported is True
    assert result.backend == "text"
    assert result.text == "ticker,price\nAAPL,200\n"
    assert result.truncated is False


def test_html_extraction_has_backend_label(tmp_path: Path) -> None:
    html_content = b"<html><body><p>Hello</p></body></html>"
    path = tmp_path / "page.html"
    path.write_bytes(html_content)
    result = _extract_text(
        path, content=html_content, mime_type="text/html", max_chars=10_000
    )
    assert result.supported is True
    assert result.backend == "html"
    assert "Hello" in (result.text or "")


# ---------------------------------------------------------------------------
# Bounding: large Office doc is truncated at max_chars
# ---------------------------------------------------------------------------


def test_docx_truncated_at_max_chars(tmp_path: Path) -> None:
    long_paragraph = "A" * 500
    path = make_docx(tmp_path, paragraph=long_paragraph, table_rows=[])
    result = _extract_docx_text(path, max_chars=100)
    assert result.supported is True
    assert result.truncated is True
    assert len(result.text or "") == 100


def test_xlsx_truncated_at_max_chars(tmp_path: Path) -> None:
    rows = [["X" * 100] * 3 for _ in range(10)]
    path = make_xlsx(tmp_path, sheet_name="BigSheet", rows=rows)
    result = _extract_xlsx_text(path, max_chars=50)
    assert result.supported is True
    assert result.truncated is True
    assert len(result.text or "") == 50


def test_pptx_truncated_at_max_chars(tmp_path: Path) -> None:
    path = make_pptx(tmp_path, title="T" * 200, body="B" * 300)
    result = _extract_pptx_text(path, max_chars=30)
    assert result.supported is True
    assert result.truncated is True
    assert len(result.text or "") == 30


# ---------------------------------------------------------------------------
# TextExtraction defaults
# ---------------------------------------------------------------------------


def test_text_extraction_new_fields_have_defaults() -> None:
    t = TextExtraction(supported=True, text="hello")
    assert t.backend == "local"
    assert t.warnings == ()
    assert t.recoverable is False


def test_text_extraction_unsupported_defaults() -> None:
    t = TextExtraction(supported=False)
    assert t.text is None
    assert t.truncated is False
    assert t.backend == "local"
    assert t.recoverable is False
