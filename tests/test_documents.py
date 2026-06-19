"""Tests for the Document Studio core (HIG-244 Phase 1).

Writer tests are pure string assertions (no DB, no binary). The render test
exercises the real ``typst`` binary and skips when it is unavailable.
"""

from __future__ import annotations

import io

import pytest
from pypdf import PdfReader

from kortny.documents import (
    PITCH_THEME,
    REPORT_THEME,
    Chart,
    DocKind,
    DocumentSpec,
    compile_chart_spec,
    render_chart_png,
    render_chart_svg,
    render_document,
    render_spec_pdf,
    resolve_theme,
    theme_names,
    typst_available,
)
from kortny.documents.typst_writer import esc, esc_str

_FULL_SPEC = {
    "doc_kind": "report",
    "title": "Quarterly $15B Review",
    "blocks": [
        {
            "type": "cover_header",
            "eyebrow": "Market Brief",
            "title": "The SpaceX IPO",
            "subtitle": "What a public offering means",
            "accent_tail": "IPO",
            "meta": ["Kortny", "2026"],
        },
        {"type": "heading", "text": "Why now"},
        {"type": "prose", "text": "First paragraph.\n\nSecond paragraph."},
        {
            "type": "stat_cards",
            "cards": [
                {"value": "$350B", "label": "Valuation", "note": "midpoint"},
                {"value": "62%", "label": "Starlink"},
            ],
        },
        {
            "type": "table",
            "caption": "Comparables",
            "columns": ["Company", "Raise", "Sector"],
            "rows": [["Arm", "$4.9B", "Semis"], ["Reddit", "$748M"]],
        },
        {"type": "callout", "label": "Takeaway", "text": "It matters."},
        {"type": "pull_quote", "text": "A public currency.", "attribution": "Analyst"},
        {"type": "cta", "label": "View more →", "text": "All the data."},
        {
            "type": "section_divider",
            "index": "01",
            "label": "Section One",
            "title": "The Setup",
            "subtitle": "How we got here",
        },
    ],
}


def _spec(**overrides: object) -> DocumentSpec:
    return DocumentSpec.model_validate({**_FULL_SPEC, **overrides})


# --------------------------------------------------------------------------- #
# Escaping
# --------------------------------------------------------------------------- #


def test_esc_escapes_math_dollar() -> None:
    # $ opens math mode in Typst and would break compilation on currency.
    assert esc("$15B") == "\\$15B"


def test_esc_escapes_all_sigils() -> None:
    assert esc("a#b[c]*_`@<>") == "a\\#b\\[c\\]\\*\\_\\`\\@\\<\\>"


def test_esc_escapes_backslash_first() -> None:
    assert esc("a\\b") == "a\\\\b"


def test_esc_str_escapes_quote_and_backslash() -> None:
    # Inside a Typst string literal only \ and " are special — NOT the markup
    # sigils. A bare " would terminate the literal and break compilation.
    assert esc_str('CEO "Q2" Brief') == 'CEO \\"Q2\\" Brief'
    assert esc_str("a\\b") == "a\\\\b"


def test_esc_str_leaves_markup_sigils_literal() -> None:
    # # and $ are literal text inside a string; markup-escaping them ("\#"/"\$")
    # would emit an invalid string escape sequence and fail to compile.
    assert esc_str("#1 grew $15B") == "#1 grew $15B"


def test_string_literal_contexts_escape_quotes_not_markup() -> None:
    # A quoted title/heading/caption must not leak a bare " into the title:,
    # #upper("…") and #eyebrow("…") string literals.
    src = render_document(
        _spec(
            title='CEO "Q2" Brief',
            blocks=[
                {"type": "heading", "text": 'The "Big" One'},
                {
                    "type": "table",
                    "caption": 'Q2 "final"',
                    "columns": ['Co "X"'],
                    "rows": [["a"]],
                },
            ],
        )
    )
    assert '#set document(title: "CEO \\"Q2\\" Brief")' in src
    assert '#upper("The \\"Big\\" One")' in src
    # No bare unescaped quote survives inside a string-literal interpolation.
    assert 'title: "CEO "Q2"' not in src


# --------------------------------------------------------------------------- #
# Theme resolution
# --------------------------------------------------------------------------- #


def test_theme_names_includes_builtins() -> None:
    names = theme_names()
    assert "report" in names
    assert "pitch" in names


def test_resolve_theme_explicit_name_wins() -> None:
    assert resolve_theme(doc_kind=DocKind.report, name="pitch") is PITCH_THEME


def test_resolve_theme_defaults_by_kind() -> None:
    assert resolve_theme(doc_kind=DocKind.pitch, name=None) is PITCH_THEME
    assert resolve_theme(doc_kind=DocKind.report, name=None) is REPORT_THEME
    # brief/memo fall back to the report theme.
    assert resolve_theme(doc_kind=DocKind.memo, name=None) is REPORT_THEME


def test_resolve_theme_unknown_name_falls_back_not_raises() -> None:
    assert resolve_theme(doc_kind=DocKind.report, name="nope") is REPORT_THEME


# --------------------------------------------------------------------------- #
# Writer output
# --------------------------------------------------------------------------- #


def test_render_document_is_deterministic() -> None:
    assert render_document(_spec()) == render_document(_spec())


def test_render_document_emits_theme_tokens() -> None:
    src = render_document(_spec(theme="report"))
    assert f'#let accent = rgb("{REPORT_THEME.colors.accent}")' in src
    assert f'#let display = "{REPORT_THEME.display_font}"' in src
    # The title is a Typst string literal, where $ is literal text (not math
    # mode), so it must NOT be markup-escaped to \$ — that would be an invalid
    # string escape and break compilation.
    assert 'set document(title: "Quarterly $15B Review")' in src


def test_cover_uses_full_bleed_dark_page_and_accent_tail() -> None:
    src = render_document(_spec())
    # Cover is a dedicated dark page, not a background div.
    assert "#page(fill: dividerBg, header: none, footer: none" in src
    # accent_tail splits the title so "IPO" is accent-coloured.
    assert "The SpaceX #text(fill: accent)[IPO]" in src


def test_section_divider_is_its_own_page() -> None:
    src = render_document(_spec())
    assert src.count("#page(fill: dividerBg") == 2  # cover + one divider


def test_prose_splits_paragraphs() -> None:
    src = render_document(_spec())
    assert "First paragraph." in src
    assert "Second paragraph." in src


def test_stat_cards_grid_matches_card_count() -> None:
    src = render_document(_spec())
    # Two cards -> a two-column grid.
    assert "#grid(columns: (1fr, 1fr), gutter: 12pt," in src


def test_table_pads_ragged_rows() -> None:
    # The second row has 2 cells for a 3-column table; it must be padded so the
    # Typst grid stays valid rather than throwing.
    src = render_document(_spec())
    # 3 columns x (1 header-less) + body; ensure Reddit row rendered 3 cells.
    assert src.count("[#text(font: bodyFont, size: 10pt)[Reddit]]") == 1
    assert "[#text(font: bodyFont, size: 10pt)[Semis]]" in src


def test_callout_panel_is_content_sized_not_fixed_height() -> None:
    src = render_document(_spec())
    # A Typst block with no height -> shrinks to content (no empty-void bug).
    assert "#block(width: 100%, fill: panelTint, inset: 18pt" in src
    assert "height" not in _callout_fragment(src)


def _callout_fragment(src: str) -> str:
    start = src.index("fill: panelTint")
    return src[start : start + 200]


def test_currency_in_card_value_is_escaped() -> None:
    src = render_document(_spec())
    assert "[\\$350B]" in src


# --------------------------------------------------------------------------- #
# Render (real typst binary)
# --------------------------------------------------------------------------- #

_TYPST_MISSING = not typst_available()


@pytest.mark.skipif(_TYPST_MISSING, reason="typst binary not installed")
def test_render_spec_pdf_produces_valid_pdf() -> None:
    pdf = render_spec_pdf(_spec())
    assert pdf.startswith(b"%PDF")
    reader = PdfReader(io.BytesIO(pdf))
    assert len(reader.pages) >= 1


@pytest.mark.skipif(_TYPST_MISSING, reason="typst binary not installed")
def test_render_spec_pdf_both_themes() -> None:
    for theme in theme_names():
        pdf = render_spec_pdf(_spec(theme=theme))
        assert pdf.startswith(b"%PDF")


@pytest.mark.skipif(_TYPST_MISSING, reason="typst binary not installed")
def test_render_spec_pdf_compiles_with_quotes_and_sigils() -> None:
    # Regression: ordinary quoted/sigil content in string-literal contexts
    # (title, headings, captions, eyebrows) must still compile, not crash Typst.
    pdf = render_spec_pdf(
        _spec(
            title='CEO "Q2" Brief — #1',
            blocks=[
                {
                    "type": "cover_header",
                    "eyebrow": 'Market "Brief"',
                    "title": "The IPO",
                    "meta": ['$15B "raise"'],
                },
                {"type": "heading", "text": 'The "Big" One'},
                {
                    "type": "table",
                    "caption": 'Q2 "final"',
                    "columns": ['Co "X"', "Raise"],
                    "rows": [["Arm", "$4.9B"]],
                },
                {"type": "callout", "label": 'Key "note"', "text": "Done."},
            ],
        )
    )
    assert pdf.startswith(b"%PDF")


# --------------------------------------------------------------------------- #
# PPTX / DOCX writers (pure-Python, no binary)
# --------------------------------------------------------------------------- #


def test_render_pptx_produces_valid_deck() -> None:
    from pptx import Presentation

    from kortny.documents import render_pptx

    data = render_pptx(_spec())
    assert data[:2] == b"PK"  # OOXML zip
    prs = Presentation(io.BytesIO(data))
    # cover -> title slide; section_divider -> section slide; at least one
    # content slide for the heading group => 3+ slides.
    assert len(list(prs.slides)) >= 3


def test_render_pptx_both_themes() -> None:
    from kortny.documents import render_pptx

    for theme in theme_names():
        assert render_pptx(_spec(theme=theme))[:2] == b"PK"


def test_render_docx_produces_valid_document() -> None:
    from docx import Document

    from kortny.documents import render_docx

    data = render_docx(_spec())
    assert data[:2] == b"PK"
    doc = Document(io.BytesIO(data))
    # stat_cards + table + callout each render as a table.
    assert len(doc.tables) >= 3
    assert any("SpaceX" in p.text for p in doc.paragraphs)


def test_render_docx_both_themes() -> None:
    from kortny.documents import render_docx

    for theme in theme_names():
        assert render_docx(_spec(theme=theme))[:2] == b"PK"


# --------------------------------------------------------------------------- #
# XLSX (data export)
# --------------------------------------------------------------------------- #


def _xlsx_sheet_names(data: bytes) -> list[str]:
    import re as _re
    import zipfile

    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        assert "[Content_Types].xml" in zf.namelist()
        workbook_xml = zf.read("xl/workbook.xml").decode("utf-8")
    return _re.findall(r'<sheet[^>]*name="([^"]+)"', workbook_xml)


def test_render_xlsx_produces_valid_workbook_with_data_sheets() -> None:
    from kortny.documents import render_xlsx

    data = render_xlsx(_spec())
    assert data[:2] == b"PK"  # OOXML zip
    names = _xlsx_sheet_names(data)
    # The fixture has a table + stat_cards → Summary + Table 1 + Metrics.
    assert "Summary" in names
    assert "Table 1" in names
    assert "Metrics" in names


def test_render_xlsx_chart_block_becomes_a_data_sheet() -> None:
    from kortny.documents import render_xlsx

    spec = _spec(blocks=[*_FULL_SPEC["blocks"], _BAR.model_dump()])
    names = _xlsx_sheet_names(render_xlsx(spec))
    assert "Chart 1 Data" in names


def test_render_xlsx_dedupes_multiple_table_sheets() -> None:
    from kortny.documents import render_xlsx

    two_tables = {
        "type": "table",
        "columns": ["A", "B"],
        "rows": [["1", "2"]],
    }
    spec = _spec(blocks=[two_tables, dict(two_tables)])
    names = _xlsx_sheet_names(render_xlsx(spec))
    assert "Table 1" in names
    assert "Table 2" in names


def test_xlsx_is_poor_fit_for_prose_good_for_data() -> None:
    from kortny.documents import xlsx_is_poor_fit

    prose_only = _spec(
        blocks=[
            {"type": "heading", "text": "Overview"},
            {"type": "prose", "text": "All narrative, no data."},
        ]
    )
    assert xlsx_is_poor_fit(prose_only) is True
    # The full fixture (table + stat_cards) is a fine spreadsheet.
    assert xlsx_is_poor_fit(_spec()) is False


# --------------------------------------------------------------------------- #
# Deterministic critique (lint + auto-fix + post-render validation)
# --------------------------------------------------------------------------- #


def test_critique_pads_ragged_table_rows() -> None:
    from kortny.documents.critique import critique_and_fix

    spec = _spec(
        blocks=[
            {"type": "table", "columns": ["A", "B", "C"], "rows": [["1"], ["2", "3"]]}
        ]
    )
    result = critique_and_fix(spec)
    table = result.spec.blocks[0]
    assert all(len(row) == 3 for row in table.rows)  # type: ignore[union-attr]
    assert any(
        i.code == "ragged_table_rows" and i.autofix == "applied" for i in result.issues
    )


def test_critique_drops_empty_table_and_blank_blocks() -> None:
    from kortny.documents.critique import critique_and_fix

    spec = _spec(
        blocks=[
            {"type": "heading", "text": "Keep me"},
            {"type": "prose", "text": "body"},
            {"type": "table", "columns": ["A"], "rows": []},
            {"type": "prose", "text": "   "},
        ]
    )
    result = critique_and_fix(spec)
    kinds = [b.type for b in result.spec.blocks]
    assert "table" not in kinds
    assert kinds.count("prose") == 1
    assert {"empty_table", "empty_prose"} <= {i.code for i in result.issues}


def test_critique_dedupes_columns_and_drops_bad_accent_tail() -> None:
    from kortny.documents.critique import critique_and_fix

    spec = _spec(
        blocks=[
            {"type": "cover_header", "title": "Report", "accent_tail": "Nope"},
            {"type": "table", "columns": ["X", "", "X"], "rows": [["1", "2", "3"]]},
        ]
    )
    result = critique_and_fix(spec)
    cover, table = result.spec.blocks
    assert cover.accent_tail is None  # type: ignore[union-attr]
    assert len(set(table.columns)) == 3  # type: ignore[union-attr]


def test_critique_empty_document_is_an_error() -> None:
    from kortny.documents.critique import critique_and_fix

    result = critique_and_fix(_spec(blocks=[{"type": "prose", "text": "  "}]))
    assert result.has_errors
    assert any(i.code == "empty_document" for i in result.issues)


def test_validate_render_flags_garbage_and_passes_real_files() -> None:
    from kortny.documents import render_xlsx
    from kortny.documents.critique import validate_render

    assert validate_render(b"not a pdf", "pdf")[0].severity == "error"
    assert validate_render(b"not a zip", "xlsx")[0].severity == "error"
    assert validate_render(render_xlsx(_spec()), "xlsx") == []


# --------------------------------------------------------------------------- #
# Canvas (markdown delivery)
# --------------------------------------------------------------------------- #


def test_render_pptx_paginates_long_table_without_truncation() -> None:
    from pptx import Presentation

    from kortny.documents import render_pptx

    big_table = {
        "type": "table",
        "columns": ["Title", "Note"],
        "rows": [[f"Row {i}", "some descriptive text " * 3] for i in range(80)],
    }
    spec = _spec(blocks=[{"type": "heading", "text": "The Year's Best"}, big_table])
    prs = Presentation(io.BytesIO(render_pptx(spec)))

    table_shapes = [sh for s in prs.slides for sh in s.shapes if sh.has_table]
    # The table spilled across continuation slides...
    assert len(table_shapes) >= 2
    # ...and every data row survived (no truncation).
    total_rows = sum(len(sh.table.rows) - 1 for sh in table_shapes)
    assert total_rows == 80


def test_critique_warns_on_unlabelled_chart() -> None:
    from kortny.documents.critique import critique_and_fix

    spec = _spec(
        blocks=[
            {
                "type": "chart",
                "chart_type": "bar",
                "series": [{"name": "s", "points": [{"x": "a", "y": 1.0}]}],
            }
        ]
    )
    codes = {i.code for i in critique_and_fix(spec).issues}
    assert "chart_missing_title" in codes
    assert "chart_missing_axis_labels" in codes


def test_render_canvas_markdown_maps_blocks_and_drops_charts() -> None:
    from kortny.documents.canvas_writer import render_canvas_markdown

    spec = _spec(blocks=[*_FULL_SPEC["blocks"], _BAR.model_dump()])
    markdown, omitted = render_canvas_markdown(spec)
    assert markdown.startswith("# ")  # cover title as H1
    assert "| " in markdown  # the table became a markdown table
    # The chart is dropped with a note.
    assert omitted
    assert "Charts aren't supported" in markdown


# --------------------------------------------------------------------------- #
# Charts (vl-convert)
# --------------------------------------------------------------------------- #

_BAR = Chart.model_validate(
    {
        "chart_type": "bar",
        "title": "Sales",
        "x_label": "Q",
        "y_label": "USD",
        "series": [
            {"name": "Rev", "points": [{"x": "Q1", "y": 3}, {"x": "Q2", "y": 5}]}
        ],
    }
)
_MULTI_LINE = Chart.model_validate(
    {
        "chart_type": "line",
        "series": [
            {"name": "A", "points": [{"x": 2024, "y": 3}, {"x": 2025, "y": 7}]},
            {"name": "B", "points": [{"x": 2024, "y": 5}, {"x": 2025, "y": 4}]},
        ],
    }
)


def test_chart_spec_single_series_uses_accent_no_legend() -> None:
    spec = compile_chart_spec(_BAR, REPORT_THEME)
    # Single series -> mark painted with the brand accent, no color encoding.
    assert spec["mark"] == {"type": "bar", "color": REPORT_THEME.colors.accent}
    assert "color" not in spec["encoding"]
    # Author category order preserved (no alphabetical re-sort).
    assert spec["encoding"]["x"]["sort"] is None


def test_chart_spec_multi_series_uses_palette_and_color() -> None:
    spec = compile_chart_spec(_MULTI_LINE, REPORT_THEME)
    assert spec["mark"] == "line"
    assert spec["encoding"]["color"]["field"] == "series"
    # Numeric x detected as quantitative.
    assert spec["encoding"]["x"]["type"] == "quantitative"
    # Brand accent leads the categorical palette; theme font applied.
    assert spec["config"]["range"]["category"][0] == REPORT_THEME.colors.accent
    assert spec["config"]["font"] == REPORT_THEME.body_font


def test_chart_spec_pie_uses_arc_theta() -> None:
    pie = Chart.model_validate(
        {"chart_type": "pie", "series": [{"name": "s", "points": [{"x": "A", "y": 1}]}]}
    )
    spec = compile_chart_spec(pie, REPORT_THEME)
    assert spec["mark"]["type"] == "arc"
    assert "theta" in spec["encoding"]


def test_render_chart_svg_and_png() -> None:
    svg = render_chart_svg(_BAR, REPORT_THEME)
    assert svg.lstrip().startswith("<svg")
    png = render_chart_png(_BAR, REPORT_THEME)
    assert png[:4] == b"\x89PNG"


def test_typst_chart_block_emits_image_and_asset() -> None:
    from kortny.documents.typst_writer import build_typst

    spec = _spec(blocks=[*_FULL_SPEC["blocks"], _BAR.model_dump()])
    source, assets = build_typst(spec)
    # The chart is the 9th block (index 8) -> chart_8.svg.
    assert any(name.endswith(".svg") for name in assets)
    asset_name = next(n for n in assets if n.endswith(".svg"))
    assert f'#image("{asset_name}"' in source
    assert assets[asset_name].lstrip().startswith(b"<svg")


@pytest.mark.skipif(_TYPST_MISSING, reason="typst binary not installed")
def test_render_spec_pdf_with_chart() -> None:
    spec = _spec(blocks=[*_FULL_SPEC["blocks"], _BAR.model_dump()])
    pdf = render_spec_pdf(spec)
    assert pdf.startswith(b"%PDF")


def test_chart_embeds_in_pptx_and_docx() -> None:
    from docx import Document as DocxDocument
    from pptx import Presentation

    from kortny.documents import render_docx, render_pptx

    spec = _spec(blocks=[*_FULL_SPEC["blocks"], _BAR.model_dump()])
    prs = Presentation(io.BytesIO(render_pptx(spec)))
    # At least one picture shape across the deck.
    assert any(
        shape.shape_type is not None and "PICTURE" in str(shape.shape_type)
        for slide in prs.slides
        for shape in slide.shapes
    )
    doc = DocxDocument(io.BytesIO(render_docx(spec)))
    assert "graphicData" in doc.element.xml  # an embedded drawing
