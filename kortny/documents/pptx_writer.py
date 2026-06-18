"""Render a Document Studio IR to a PowerPoint deck (HIG-244 Phase 2).

The same format-agnostic IR that drives the Typst PDF writer drives this one —
the agent emits one JSON spec and picks a format. A page-oriented document maps
onto slides heuristically:

* ``cover_header`` -> a title slide (dark, matching the PDF cover).
* ``section_divider`` -> a section slide.
* ``heading`` -> starts a new content slide whose title is the heading; the
  blocks that follow (prose, stat_cards, table, callout, pull_quote, cta) flow
  onto that slide until the next heading/divider. A slide that fills up spills
  onto a continuation slide rather than overflowing off the canvas.

Built with python-pptx (pure-Python, native-editable .pptx, no browser). Theme
colours/fonts come from the same Theme object the Typst writer uses.
"""

from __future__ import annotations

import io
from typing import cast

from pptx import Presentation as new_presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.text.text import _Paragraph
from pptx.util import Emu, Inches, Length, Pt

from kortny.documents.ir import (
    CTA,
    Callout,
    CoverHeader,
    DocumentSpec,
    Heading,
    Prose,
    PullQuote,
    SectionDivider,
    StatCards,
    Table,
)
from kortny.documents.themes import Theme, resolve_theme

# 16:9 canvas.
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)
_MARGIN = Inches(0.9)
_CONTENT_TOP = Inches(1.7)
_CONTENT_BOTTOM = Inches(7.0)
# Length arithmetic is typed as int in the stubs, so rebuild as Emu (a Length).
_BODY_W: Length = Emu(int(_SLIDE_W) - 2 * int(_MARGIN))


def _rgb(hex_color: str) -> RGBColor:
    return cast(RGBColor, RGBColor.from_string(hex_color.lstrip("#")))


def render_pptx(spec: DocumentSpec) -> bytes:
    """Render ``spec`` to a .pptx file as bytes."""

    theme = resolve_theme(doc_kind=spec.doc_kind, name=spec.theme)
    prs = new_presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H
    builder = _DeckBuilder(prs, theme)
    builder.build(spec)
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


class _DeckBuilder:
    def __init__(self, prs: Presentation, theme: Theme) -> None:
        self.prs = prs
        self.theme = theme
        self.blank = prs.slide_layouts[6]  # blank layout
        self._content: Slide | None = None
        self._y: int = int(_CONTENT_TOP)
        self._content_title: str | None = None

    # -- slide lifecycle --------------------------------------------------- #

    def _new_slide(self, *, dark: bool) -> Slide:
        slide = self.prs.slides.add_slide(self.blank)
        fill = slide.background.fill
        fill.solid()
        c = self.theme.colors
        fill.fore_color.rgb = _rgb(c.divider_bg if dark else c.paper)
        return cast(Slide, slide)

    def _ensure_content_slide(self, *, title: str | None) -> None:
        """Open a fresh content slide titled ``title``."""

        self._content = self._new_slide(dark=False)
        self._content_title = title
        self._y = int(_CONTENT_TOP)
        if title:
            self._add_text(
                self._content,
                title,
                left=_MARGIN,
                top=Inches(0.7),
                width=_BODY_W,
                height=Inches(0.9),
                size=26,
                bold=True,
                font=self.theme.display_font,
                color=self.theme.colors.ink,
            )
            # accent rule under the title
            rule = self._content.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, _MARGIN, Inches(1.5), Inches(0.5), Pt(3)
            )
            rule.fill.solid()
            rule.fill.fore_color.rgb = _rgb(self.theme.colors.accent)
            rule.line.fill.background()
            _no_shadow(rule)

    def _space(self, needed: int) -> None:
        """Ensure ``needed`` EMU of vertical room, spilling to a new slide."""

        if self._content is None or self._y + needed > int(_CONTENT_BOTTOM):
            title = self._content_title
            self._ensure_content_slide(title=f"{title} (cont.)" if title else None)

    # -- public build ------------------------------------------------------ #

    def build(self, spec: DocumentSpec) -> None:
        for block in spec.blocks:
            self._render(block)

    def _render(self, block: object) -> None:
        if isinstance(block, CoverHeader):
            self._title_slide(block)
        elif isinstance(block, SectionDivider):
            self._section_slide(block)
        elif isinstance(block, Heading):
            self._ensure_content_slide(title=block.text)
        elif isinstance(block, Prose):
            self._prose(block)
        elif isinstance(block, StatCards):
            self._stat_cards(block)
        elif isinstance(block, Table):
            self._table(block)
        elif isinstance(block, Callout):
            self._callout(block)
        elif isinstance(block, PullQuote):
            self._pull_quote(block)
        elif isinstance(block, CTA):
            self._cta(block)

    # -- title / section --------------------------------------------------- #

    def _title_slide(self, block: CoverHeader) -> None:
        c = self.theme.colors
        slide = self._new_slide(dark=self.theme.dark_dividers)
        fg = c.on_dark if self.theme.dark_dividers else c.ink
        fg_muted = c.on_dark_muted if self.theme.dark_dividers else c.muted
        # Stack eyebrow + title + subtitle in one anchored textbox so spacing is
        # natural and a short (1-line) title never leaves a dead gap.
        box = slide.shapes.add_textbox(_MARGIN, Inches(2.3), _BODY_W, Inches(3.6))
        tf = box.text_frame
        tf.word_wrap = True
        first = True

        def _para() -> _Paragraph:
            nonlocal first
            if first:
                first = False
                return cast(_Paragraph, tf.paragraphs[0])
            return cast(_Paragraph, tf.add_paragraph())

        if block.eyebrow:
            p = _para()
            _set_run(
                p,
                block.eyebrow.upper(),
                size=12,
                font=self.theme.mono_font,
                color=c.accent,
            )
            p.space_after = Pt(12)
        title_p = _para()
        self._title_runs(title_p, block.title, block.accent_tail, fg)
        title_p.space_after = Pt(16)
        if block.subtitle:
            p = _para()
            _set_run(
                p, block.subtitle, size=18, font=self.theme.body_font, color=fg_muted
            )
        if block.meta:
            self._add_text(
                slide,
                "   ·   ".join(block.meta).upper(),
                left=_MARGIN,
                top=Inches(6.7),
                width=_BODY_W,
                height=Inches(0.4),
                size=10,
                font=self.theme.mono_font,
                color=fg_muted,
            )
        # reset content flow so the next heading opens a fresh slide
        self._content = None

    def _title_runs(
        self, paragraph: _Paragraph, title: str, accent_tail: str | None, fg: str
    ) -> None:
        """Big display title, optionally accent-colouring a trailing fragment."""

        if accent_tail and accent_tail != title and title.endswith(accent_tail):
            head = title[: -len(accent_tail)]
            _set_run(
                paragraph,
                head,
                size=46,
                bold=True,
                font=self.theme.display_font,
                color=fg,
            )
            _set_run(
                paragraph,
                accent_tail,
                size=46,
                bold=True,
                font=self.theme.display_font,
                color=self.theme.colors.accent,
            )
        else:
            _set_run(
                paragraph,
                title,
                size=46,
                bold=True,
                font=self.theme.display_font,
                color=fg,
            )

    def _section_slide(self, block: SectionDivider) -> None:
        c = self.theme.colors
        dark = self.theme.dark_dividers
        slide = self._new_slide(dark=dark)
        fg = c.on_dark if dark else c.ink
        fg_muted = c.on_dark_muted if dark else c.muted
        box = slide.shapes.add_textbox(_MARGIN, Inches(2.7), _BODY_W, Inches(2.6))
        tf = box.text_frame
        tf.word_wrap = True
        first = True

        def _para() -> _Paragraph:
            nonlocal first
            if first:
                first = False
                return cast(_Paragraph, tf.paragraphs[0])
            return cast(_Paragraph, tf.add_paragraph())

        if block.index:
            p = _para()
            _set_run(p, block.index, size=20, font=self.theme.mono_font, color=c.accent)
            p.space_after = Pt(8)
        if block.label:
            p = _para()
            _set_run(
                p,
                block.label.upper(),
                size=12,
                font=self.theme.mono_font,
                color=fg_muted,
            )
            p.space_after = Pt(10)
        title_p = _para()
        _set_run(
            title_p,
            block.title,
            size=34,
            bold=True,
            font=self.theme.display_font,
            color=fg,
        )
        if block.subtitle:
            title_p.space_after = Pt(14)
            p = _para()
            _set_run(
                p, block.subtitle, size=14, font=self.theme.body_font, color=fg_muted
            )
        self._content = None

    # -- content blocks ---------------------------------------------------- #

    def _prose(self, block: Prose) -> None:
        text = block.text.strip()
        height = _estimate_text_height(text, chars_per_line=95, line_h=0.28)
        self._space(height + int(Inches(0.15)))
        assert self._content is not None
        self._add_text(
            self._content,
            text,
            left=_MARGIN,
            top=Emu(self._y),
            width=_BODY_W,
            height=Emu(height),
            size=14,
            font=self.theme.body_font,
            color=self.theme.colors.ink,
        )
        self._y += height + int(Inches(0.2))

    def _stat_cards(self, block: StatCards) -> None:
        card_h = int(Inches(1.4))
        self._space(card_h + int(Inches(0.2)))
        assert self._content is not None
        n = len(block.cards)
        gutter = int(Inches(0.2))
        card_w = (int(_BODY_W) - gutter * (n - 1)) // n
        c = self.theme.colors
        for i, card in enumerate(block.cards):
            left = int(_MARGIN) + i * (card_w + gutter)
            box = self._content.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(left),
                Emu(self._y),
                Emu(card_w),
                Emu(card_h),
            )
            box.fill.solid()
            box.fill.fore_color.rgb = _rgb(c.card)
            box.line.fill.background()
            _no_shadow(box)
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.15)
            tf.margin_top = Inches(0.1)
            tf.margin_bottom = Inches(0.1)
            _set_run(
                tf.paragraphs[0],
                card.value,
                size=26,
                bold=True,
                font=self.theme.display_font,
                color=c.accent,
            )
            p_label = tf.add_paragraph()
            _set_run(
                p_label,
                card.label,
                size=12,
                bold=True,
                font=self.theme.body_font,
                color=c.ink,
            )
            if card.note:
                p_note = tf.add_paragraph()
                _set_run(
                    p_note, card.note, size=9, font=self.theme.mono_font, color=c.muted
                )
        self._y += card_h + int(Inches(0.25))

    def _table(self, block: Table) -> None:
        rows = len(block.rows) + 1
        cols = len(block.columns)
        row_h = int(Inches(0.4))
        total_h = row_h * rows
        self._space(total_h + int(Inches(0.2)))
        assert self._content is not None
        c = self.theme.colors
        shape = self._content.shapes.add_table(
            rows, cols, _MARGIN, Emu(self._y), _BODY_W, Emu(total_h)
        )
        table = shape.table
        # The default table style paints a banded blue Office theme that fights
        # our explicit cell fills; turn the styled bands off.
        table.first_row = False
        table.horz_banding = False
        for j, col in enumerate(block.columns):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(c.ink)
            _set_run(
                cell.text_frame.paragraphs[0],
                col.upper(),
                size=10,
                bold=True,
                font=self.theme.mono_font,
                color="#FFFFFF",
            )
        for i, row in enumerate(block.rows, start=1):
            padded = (list(row) + [""] * cols)[:cols]
            for j, value in enumerate(padded):
                cell = table.cell(i, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(c.paper)
                _set_run(
                    cell.text_frame.paragraphs[0],
                    str(value),
                    size=11,
                    font=self.theme.body_font,
                    color=c.ink,
                )
        self._y += total_h + int(Inches(0.25))

    def _callout(self, block: Callout) -> None:
        text = block.text.strip()
        label_h = int(Inches(0.32)) if block.label else 0
        body_h = _estimate_text_height(text, chars_per_line=80, line_h=0.3)
        total = body_h + label_h + int(Inches(0.4))
        self._space(total + int(Inches(0.2)))
        assert self._content is not None
        c = self.theme.colors
        box = self._content.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, _MARGIN, Emu(self._y), _BODY_W, Emu(total)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(self.theme.panel_tint)
        box.line.color.rgb = _rgb(c.accent)
        _no_shadow(box)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.15)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        if block.label:
            _set_run(
                tf.paragraphs[0],
                block.label.upper(),
                size=11,
                bold=True,
                font=self.theme.mono_font,
                color=c.accent,
            )
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        _set_run(p, text, size=14, font=self.theme.body_font, color=c.ink)
        self._y += total + int(Inches(0.2))

    def _pull_quote(self, block: PullQuote) -> None:
        text = f'"{block.text.strip()}"'
        body_h = _estimate_text_height(text, chars_per_line=70, line_h=0.4)
        total = body_h + int(Inches(0.5))
        self._space(total + int(Inches(0.2)))
        assert self._content is not None
        c = self.theme.colors
        bar = self._content.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, _MARGIN, Emu(self._y), Pt(3), Emu(total)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(c.accent)
        bar.line.fill.background()
        _no_shadow(bar)
        tf_left = int(_MARGIN) + int(Inches(0.25))
        box = self._content.shapes.add_textbox(
            Emu(tf_left),
            Emu(self._y),
            Emu(int(_BODY_W) - int(Inches(0.25))),
            Emu(total),
        )
        tf = box.text_frame
        tf.word_wrap = True
        _set_run(
            tf.paragraphs[0],
            text,
            size=18,
            italic=True,
            font=self.theme.display_font,
            color=c.ink,
        )
        if block.attribution:
            p = tf.add_paragraph()
            _set_run(
                p,
                f"— {block.attribution.upper()}",
                size=10,
                font=self.theme.mono_font,
                color=c.muted,
            )
        self._y += total + int(Inches(0.2))

    def _cta(self, block: CTA) -> None:
        self._space(int(Inches(0.8)))
        assert self._content is not None
        c = self.theme.colors
        if block.text:
            self._add_text(
                self._content,
                block.text,
                left=_MARGIN,
                top=Emu(self._y),
                width=_BODY_W,
                height=Inches(0.4),
                size=13,
                font=self.theme.body_font,
                color=c.muted,
            )
            self._y += int(Inches(0.45))
        pill = self._content.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            _MARGIN,
            Emu(self._y),
            Inches(3.2),
            Inches(0.5),
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = _rgb(c.ink)
        pill.line.fill.background()
        _no_shadow(pill)
        pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_run(
            pill.text_frame.paragraphs[0],
            block.label.upper(),
            size=11,
            bold=True,
            font=self.theme.mono_font,
            color="#FFFFFF",
            align=PP_ALIGN.CENTER,
        )
        self._y += int(Inches(0.7))

    # -- helpers ----------------------------------------------------------- #

    def _add_text(
        self,
        slide: Slide,
        text: str,
        *,
        left: int,
        top: int,
        width: int,
        height: int,
        size: int,
        font: str,
        color: str,
        bold: bool = False,
        italic: bool = False,
        align: PP_ALIGN = PP_ALIGN.LEFT,
    ) -> None:
        box = slide.shapes.add_textbox(
            Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(int(height))
        )
        tf = box.text_frame
        tf.word_wrap = True
        _set_run(
            tf.paragraphs[0],
            text,
            size=size,
            bold=bold,
            italic=italic,
            font=font,
            color=color,
            align=align,
        )


def _set_run(
    paragraph: _Paragraph,
    text: str,
    *,
    size: int,
    font: str,
    color: str,
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    # Auto-shapes default to centred text; force alignment so card/callout/pill
    # copy reads left like the PDF rather than centred-and-cramped.
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = _rgb(color)


def _no_shadow(shape: object) -> None:
    """Drop the default auto-shape drop shadow (looks dated/heavy).

    ``shadow.inherit = False`` adds an empty effect list, but the shape's
    ``<p:style>`` still carries a theme ``effectRef`` that renderers honour, so
    we also strip the style element to fully clear the inherited shadow.
    """

    shape.shadow.inherit = False  # type: ignore[attr-defined]
    element = shape._element  # type: ignore[attr-defined]
    style = element.find(qn("p:style"))
    if style is not None:
        element.remove(style)


def _estimate_text_height(text: str, *, chars_per_line: int, line_h: float) -> int:
    """Rough EMU height for wrapped text, used to flow blocks down a slide."""

    lines = 0
    for paragraph in text.split("\n"):
        lines += max(1, -(-len(paragraph) // chars_per_line))  # ceil division
    return int(Inches(max(0.3, lines * line_h)))
