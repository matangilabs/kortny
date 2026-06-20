"""Slack file download and text extraction tool."""

from __future__ import annotations

import hashlib
import logging
import mimetypes
import re
import zipfile
from collections.abc import Callable, Mapping, Sequence
from html import unescape
from html.parser import HTMLParser
from io import BytesIO
from pathlib import Path
from typing import Any, Protocol
from urllib.parse import unquote, urlparse

import httpx
from pypdf import PdfReader
from slack_sdk.errors import SlackApiError
from sqlalchemy.orm import Session

from kortny.pdf_raster import rasterize_pdf_pages
from kortny.tools.types import JsonObject, JsonSchema, ToolResult

logger = logging.getLogger(__name__)

DEFAULT_MAX_FILE_SIZE_BYTES = 25 * 1024 * 1024
DEFAULT_EXTRACTED_TEXT_MAX_CHARS = 100_000
SAFE_FILENAME_RE = re.compile(r"[^A-Za-z0-9._-]+")
SLACK_FILE_ID_RE = re.compile(r"^F[A-Z0-9]+$")

# Minimum character count to consider a PDF "text-bearing".  Scanned PDFs
# typically extract 0 characters; native PDFs may have a handful of whitespace
# characters; anything below this threshold triggers OCR when available.
SCANNED_PDF_TEXT_THRESHOLD = 16

# A callable that takes a sequence of page-image PNG bytes and returns
# transcribed Markdown.  Injected into SlackFileReadTool so the tool stays
# testable without a real LLM.
PdfPageOcr = Callable[[Sequence[bytes]], str]

# Zip-bomb guard constants for Office files (.docx/.xlsx/.pptx are ZIP archives)
MAX_OFFICE_UNCOMPRESSED_BYTES = 200 * 1024 * 1024  # 200 MB
MAX_OFFICE_ZIP_ENTRIES = 5000
MAX_OFFICE_COMPRESSION_RATIO = 200


class SlackFileReadError(RuntimeError):
    """Raised when a Slack file cannot be read."""


class SlackFileLookupError(SlackFileReadError):
    """Raised when Slack file metadata lookup fails."""

    def __init__(self, code: str, message: str) -> None:
        super().__init__(message)
        self.code = code


class SlackFileReadClient(Protocol):
    """Subset of Slack WebClient used for file metadata retrieval."""

    def files_info(self, *, file: str) -> Any:
        """Fetch Slack file metadata."""


class SlackFileReadTool:
    """Download a Slack file and extract text when the type is supported."""

    name = "slack_file_read"
    description = (
        "Downloads a Slack file by file_id or private Slack file URL, saves it in "
        "the current task workspace, and extracts text for PDFs, plain text, "
        "Markdown, CSV, HTML, DOCX, XLSX, and PPTX."
    )
    parameters: JsonSchema = {
        "type": "object",
        "properties": {
            "file_id": {
                "type": "string",
                "description": "Slack file ID, such as F123ABC.",
            },
            "file_url": {
                "type": "string",
                "description": "Private Slack file URL to download with the bot token.",
            },
        },
        "additionalProperties": False,
    }

    def __init__(
        self,
        *,
        client: SlackFileReadClient,
        bot_token: str,
        working_dir: str | Path,
        max_file_size_bytes: int = DEFAULT_MAX_FILE_SIZE_BYTES,
        extracted_text_max_chars: int = DEFAULT_EXTRACTED_TEXT_MAX_CHARS,
        timeout: float = 30.0,
        transport: httpx.BaseTransport | None = None,
        session: Session | None = None,
        pdf_ocr: PdfPageOcr | None = None,
        pdf_ocr_max_pages: int = 20,
    ) -> None:
        if not bot_token.strip():
            raise ValueError("SLACK_BOT_TOKEN is required for slack_file_read")
        if max_file_size_bytes < 1:
            raise ValueError("max_file_size_bytes must be at least 1")
        if extracted_text_max_chars < 1:
            raise ValueError("extracted_text_max_chars must be at least 1")

        self.client = client
        self.bot_token = bot_token
        self.working_dir = Path(working_dir)
        self.max_file_size_bytes = max_file_size_bytes
        self.extracted_text_max_chars = extracted_text_max_chars
        self.timeout = timeout
        self.transport = transport
        self._session: Session | None = session
        self._pdf_ocr: PdfPageOcr | None = pdf_ocr
        self._pdf_ocr_max_pages: int = pdf_ocr_max_pages

    def invoke(self, args: JsonObject) -> ToolResult:
        """Download the file, extract supported text, and return metadata."""

        request = _file_request(args)
        if request.file_id is not None and not SLACK_FILE_ID_RE.match(request.file_id):
            return _recoverable_lookup_result(
                request,
                code="invalid_file_id",
                message=(
                    "slack_file_read file_id must be a Slack file ID like F123ABC, "
                    "not a message timestamp or filename"
                ),
            )

        try:
            metadata = self._metadata(request)
        except SlackFileLookupError as exc:
            return _recoverable_lookup_result(
                request,
                code=exc.code,
                message=str(exc),
            )
        if metadata.size_bytes is not None:
            _ensure_size_allowed(metadata.size_bytes, self.max_file_size_bytes)

        content, response_mime_type = self._download(metadata.download_url)
        _ensure_size_allowed(len(content), self.max_file_size_bytes)

        filename = _safe_filename(metadata.filename)
        mime_type = _mime_type(
            metadata.mime_type,
            response_mime_type,
            filename,
        )
        self.working_dir.mkdir(parents=True, exist_ok=True)
        output_path = self.working_dir / filename
        output_path.write_bytes(content)

        content_sha256 = hashlib.sha256(content).hexdigest()
        extraction_cache_status: str | None = None
        extraction: TextExtraction | None = None

        if self._session is not None:
            from kortny.tools.file_extraction_cache import FileExtractionCacheRepository

            repo = FileExtractionCacheRepository(self._session)
            extraction = repo.get(content_sha256)
            if extraction is not None:
                extraction_cache_status = "hit"

        if extraction is None:
            extraction = _extract_text(
                output_path,
                content=content,
                mime_type=mime_type,
                max_chars=self.extracted_text_max_chars,
            )
            # Scanned-PDF OCR path (HIG-279 slice 3b-2).
            # A PDF whose text layer is empty/near-empty is likely a scan.
            # If a vision-OCR callable is wired in, rasterize pages and
            # transcribe them; the OCR result is what gets cached so
            # subsequent reads are free (no re-rasterisation).
            if _is_pdf(output_path, mime_type) and _is_scanned_pdf(extraction):
                extraction = self._ocr_scanned_pdf(
                    output_path,
                    max_chars=self.extracted_text_max_chars,
                )
            if self._session is not None:
                from kortny.tools.file_extraction_cache import (
                    FileExtractionCacheRepository,
                )

                repo = FileExtractionCacheRepository(self._session)
                repo.put(content_sha256, extraction, byte_size=len(content))
                extraction_cache_status = "miss"

        output: JsonObject = {
            "file_id": metadata.file_id,
            "filename": filename,
            "mime_type": mime_type,
            "path": str(output_path),
            "size_bytes": len(content),
            "extraction_supported": extraction.supported,
            "backend": extraction.backend,
        }
        if extraction.text is not None:
            output["extracted_text"] = extraction.text
            output["extracted_text_chars"] = len(extraction.text)
            output["extracted_text_truncated"] = extraction.truncated
        if extraction.warnings:
            output["warnings"] = list(extraction.warnings)
        if extraction_cache_status is not None:
            output["extraction_cache"] = extraction_cache_status

        return ToolResult(output=output)

    def _metadata(self, request: FileRequest) -> FileMetadata:
        if request.file_url is not None:
            return FileMetadata(
                file_id=None,
                filename=_filename_from_url(request.file_url),
                mime_type=None,
                size_bytes=None,
                download_url=request.file_url,
            )

        assert request.file_id is not None
        try:
            response = self.client.files_info(file=request.file_id)
        except SlackApiError as exc:
            code = _slack_api_error_code(exc.response)
            raise SlackFileLookupError(
                code,
                f"files.info failed: {code}",
            ) from exc
        payload = _response_payload(response, "files.info")
        raw_file = payload.get("file")
        if not isinstance(raw_file, Mapping):
            raise SlackFileReadError("files.info response is missing file metadata")

        download_url = _optional_string(
            raw_file.get("url_private_download")
        ) or _optional_string(raw_file.get("url_private"))
        if download_url is None:
            raise SlackFileReadError("Slack file metadata is missing a download URL")

        return FileMetadata(
            file_id=request.file_id,
            filename=_file_metadata_filename(raw_file),
            mime_type=_optional_string(raw_file.get("mimetype")),
            size_bytes=_optional_int(raw_file.get("size")),
            download_url=download_url,
        )

    def _ocr_scanned_pdf(self, path: Path, *, max_chars: int) -> TextExtraction:
        """Rasterize a scanned PDF and transcribe it via the injected OCR callable.

        Returns a ``TextExtraction`` with ``backend="vision_ocr"`` on success,
        or a recoverable unsupported extraction when OCR is not wired in or
        rasterisation/transcription fails.  Never raises.
        """
        if self._pdf_ocr is None:
            # No vision model configured — tell the agent how to fix it.
            return TextExtraction(
                supported=False,
                recoverable=True,
                backend="pdf_textlayer",
                warnings=("scanned_pdf_needs_vision_model",),
            )

        try:
            page_pngs = _rasterize_pdf_pages(path, max_pages=self._pdf_ocr_max_pages)
        except Exception:
            logger.exception("pdf_ocr: rasterisation failed for %s", path)
            return TextExtraction(
                supported=False,
                recoverable=True,
                backend="pdf_textlayer",
                warnings=("pdf_ocr_rasterize_error",),
            )

        if not page_pngs:
            return TextExtraction(
                supported=False,
                recoverable=True,
                backend="pdf_textlayer",
                warnings=("pdf_ocr_no_pages",),
            )

        try:
            import pypdfium2  # noqa: PLC0415

            doc = pypdfium2.PdfDocument(str(path))
            total_pages = len(doc)
            doc.close()
        except Exception:
            total_pages = len(page_pngs)

        truncated_pages = total_pages > self._pdf_ocr_max_pages
        actual_page_count = min(total_pages, self._pdf_ocr_max_pages)

        try:
            ocr_text = self._pdf_ocr(page_pngs)
        except Exception:
            logger.exception("pdf_ocr: transcription failed for %s", path)
            return TextExtraction(
                supported=False,
                recoverable=True,
                backend="pdf_textlayer",
                warnings=("pdf_ocr_transcription_error",),
            )

        warnings: tuple[str, ...] = ()
        if truncated_pages:
            warnings = (f"pdf_ocr_truncated_at_{actual_page_count}_pages",)

        text = ocr_text[:max_chars]
        return TextExtraction(
            supported=True,
            text=text,
            truncated=len(ocr_text) > max_chars,
            backend="vision_ocr",
            warnings=warnings,
        )

    def _download(self, url: str) -> tuple[bytes, str | None]:
        headers = {"Authorization": f"Bearer {self.bot_token}"}
        chunks: list[bytes] = []
        total_size = 0
        with (
            httpx.Client(
                transport=self.transport,
                timeout=self.timeout,
                follow_redirects=True,
            ) as client,
            client.stream("GET", url, headers=headers) as response,
        ):
            response.raise_for_status()
            content_length = _optional_int(response.headers.get("content-length"))
            if content_length is not None:
                _ensure_size_allowed(content_length, self.max_file_size_bytes)

            for chunk in response.iter_bytes():
                total_size += len(chunk)
                _ensure_size_allowed(total_size, self.max_file_size_bytes)
                chunks.append(chunk)

            content_type = _content_type(response.headers.get("content-type"))

        return b"".join(chunks), content_type


class FileRequest:
    def __init__(self, *, file_id: str | None, file_url: str | None) -> None:
        self.file_id = file_id
        self.file_url = file_url


class FileMetadata:
    def __init__(
        self,
        *,
        file_id: str | None,
        filename: str,
        mime_type: str | None,
        size_bytes: int | None,
        download_url: str,
    ) -> None:
        self.file_id = file_id
        self.filename = filename
        self.mime_type = mime_type
        self.size_bytes = size_bytes
        self.download_url = download_url


class TextExtraction:
    def __init__(
        self,
        *,
        supported: bool,
        text: str | None = None,
        truncated: bool = False,
        backend: str = "local",
        warnings: tuple[str, ...] = (),
        recoverable: bool = False,
    ) -> None:
        self.supported = supported
        self.text = text
        self.truncated = truncated
        self.backend = backend
        self.warnings = warnings
        self.recoverable = recoverable


def _file_request(args: Mapping[str, Any]) -> FileRequest:
    file_id = _optional_string(args.get("file_id"))
    file_url = _optional_string(args.get("file_url"))
    if file_id is None and file_url is None:
        raise ValueError("slack_file_read requires 'file_id' or 'file_url'")
    if file_id is not None and file_url is not None:
        raise ValueError("slack_file_read accepts only one of 'file_id' or 'file_url'")
    return FileRequest(file_id=file_id, file_url=file_url)


def _response_payload(response: object, method: str) -> Mapping[str, Any]:
    if isinstance(response, Mapping):
        payload = response
    else:
        data = getattr(response, "data", None)
        if not isinstance(data, Mapping):
            raise SlackFileReadError(f"{method} returned a non-object response")
        payload = data

    if payload.get("ok") is False:
        error = payload.get("error")
        if not isinstance(error, str) or not error:
            error = "unknown_error"
        if method == "files.info":
            raise SlackFileLookupError(error, f"{method} failed: {error}")
        raise SlackFileReadError(f"{method} failed: {error}")

    return payload


def _recoverable_lookup_result(
    request: FileRequest,
    *,
    code: str,
    message: str,
) -> ToolResult:
    return ToolResult(
        output={
            "file_id": request.file_id,
            "file_url": request.file_url,
            "error": {
                "code": code,
                "message": message,
                "recoverable": True,
            },
        }
    )


def _slack_api_error_code(response: object) -> str:
    if isinstance(response, Mapping):
        error = response.get("error")
        if isinstance(error, str) and error:
            return error

    data = getattr(response, "data", None)
    if isinstance(data, Mapping):
        error = data.get("error")
        if isinstance(error, str) and error:
            return error

    get = getattr(response, "get", None)
    if callable(get):
        error = get("error")
        if isinstance(error, str) and error:
            return error

    return "unknown_error"


def _file_metadata_filename(raw_file: Mapping[str, Any]) -> str:
    for key in ("name", "title"):
        value = _optional_string(raw_file.get(key))
        if value is not None:
            return value
    file_id = _optional_string(raw_file.get("id"))
    if file_id is not None:
        return file_id
    return "slack-file"


def _filename_from_url(url: str) -> str:
    parsed = urlparse(url)
    name = unquote(Path(parsed.path).name)
    return name or "slack-file"


def _safe_filename(raw_filename: str) -> str:
    safe_name = SAFE_FILENAME_RE.sub("_", Path(raw_filename).name.strip())
    if safe_name in ("", ".", ".."):
        return "slack-file"
    return safe_name


def _mime_type(
    metadata_mime_type: str | None,
    response_mime_type: str | None,
    filename: str,
) -> str:
    guessed_type, _ = mimetypes.guess_type(filename)
    return (
        metadata_mime_type
        or response_mime_type
        or guessed_type
        or "application/octet-stream"
    )


def _check_office_zip_bomb(content: bytes) -> TextExtraction | None:
    """Inspect Office file ZIP structure for zip-bomb characteristics.

    Returns a TextExtraction (unsupported/recoverable) if the file is unsafe,
    or None if the file is safe to parse.
    """
    try:
        with zipfile.ZipFile(BytesIO(content)) as zf:
            entries = zf.infolist()
    except zipfile.BadZipFile:
        return TextExtraction(supported=False)

    if len(entries) > MAX_OFFICE_ZIP_ENTRIES:
        return TextExtraction(
            supported=False,
            recoverable=True,
            warnings=("office_file_exceeds_safe_limits",),
        )

    total_uncompressed = sum(e.file_size for e in entries)
    if total_uncompressed > MAX_OFFICE_UNCOMPRESSED_BYTES:
        return TextExtraction(
            supported=False,
            recoverable=True,
            warnings=("office_file_exceeds_safe_limits",),
        )

    compressed = sum(e.compress_size for e in entries)
    if (
        compressed > 0
        and total_uncompressed / compressed > MAX_OFFICE_COMPRESSION_RATIO
    ):
        return TextExtraction(
            supported=False,
            recoverable=True,
            warnings=("office_file_exceeds_safe_limits",),
        )

    return None


def _extract_text(
    path: Path,
    *,
    content: bytes,
    mime_type: str,
    max_chars: int,
) -> TextExtraction:
    if _is_pdf(path, mime_type):
        text = _extract_pdf_text(path)
        result = _bounded_text(text, max_chars=max_chars)
        result.backend = "pdf_textlayer"
        return result
    if _is_html(path, mime_type):
        text = _extract_html_text(content)
        result = _bounded_text(text, max_chars=max_chars)
        result.backend = "html"
        return result
    if _is_text_like(path, mime_type):
        text = _decode_text(content)
        result = _bounded_text(text, max_chars=max_chars)
        result.backend = "text"
        return result
    if _is_docx(path, mime_type):
        bomb = _check_office_zip_bomb(content)
        if bomb is not None:
            return bomb
        return _extract_docx_text(path, max_chars=max_chars)
    if _is_xlsx(path, mime_type):
        bomb = _check_office_zip_bomb(content)
        if bomb is not None:
            return bomb
        return _extract_xlsx_text(path, max_chars=max_chars)
    if _is_pptx(path, mime_type):
        bomb = _check_office_zip_bomb(content)
        if bomb is not None:
            return bomb
        return _extract_pptx_text(path, max_chars=max_chars)
    return TextExtraction(supported=False)


def _is_scanned_pdf(extraction: TextExtraction) -> bool:
    """Return True if the extraction looks like it came from a scanned PDF.

    A PDF whose text layer is empty or has fewer than SCANNED_PDF_TEXT_THRESHOLD
    characters after stripping whitespace is treated as scanned/image-only.
    """
    text = extraction.text or ""
    return len(text.strip()) < SCANNED_PDF_TEXT_THRESHOLD


def _rasterize_pdf_pages(path: Path, *, max_pages: int) -> list[bytes]:
    """Render up to *max_pages* pages of a PDF to PNG bytes for OCR.

    Thin wrapper over the shared ``kortny.pdf_raster.rasterize_pdf_pages`` so the
    tests' monkeypatch target stays stable; the rasterization itself lives in one
    place (also used by the Document Studio visual critic).
    """
    return rasterize_pdf_pages(path, max_pages=max_pages)


def _extract_pdf_text(path: Path) -> str:
    try:
        reader = PdfReader(str(path))
        return "\n\n".join((page.extract_text() or "").strip() for page in reader.pages)
    except Exception as exc:
        raise SlackFileReadError(f"PDF text extraction failed: {exc}") from exc


def _extract_html_text(content: bytes) -> str:
    parser = _HTMLTextExtractor()
    parser.feed(_decode_text(content))
    parser.close()
    return parser.text()


def _decode_text(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


def _extract_docx_text(path: Path, *, max_chars: int) -> TextExtraction:
    try:
        import docx
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        doc = docx.Document(str(path))
        parts: list[str] = []
        for child in doc.element.body:
            tag = child.tag
            if tag == qn("w:p"):
                para = Paragraph(child, doc)
                text = para.text
                if text:
                    parts.append(text)
            elif tag == qn("w:tbl"):
                table = Table(child, doc)
                for row in table.rows:
                    row_cells = " | ".join(cell.text for cell in row.cells)
                    if row_cells.strip():
                        parts.append(row_cells)
        text = "\n".join(parts)
        result = _bounded_text(text, max_chars=max_chars)
        result.backend = "docx"
        return result
    except Exception:
        return TextExtraction(
            supported=False,
            recoverable=True,
            warnings=("docx_parse_error",),
        )


def _extract_xlsx_text(path: Path, *, max_chars: int) -> TextExtraction:
    try:
        import openpyxl

        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"# Sheet: {sheet_name}")
            for row in ws.iter_rows():
                cells = " | ".join(
                    str(cell.value) for cell in row if cell.value is not None
                )
                if cells.strip():
                    parts.append(cells)
        wb.close()
        text = "\n".join(parts)
        result = _bounded_text(text, max_chars=max_chars)
        result.backend = "xlsx"
        return result
    except Exception:
        return TextExtraction(
            supported=False,
            recoverable=True,
            warnings=("xlsx_parse_error",),
        )


def _extract_pptx_text(path: Path, *, max_chars: int) -> TextExtraction:
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        parts: list[str] = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            parts.append(f"## Slide {slide_num}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text
                        if text.strip():
                            parts.append(text)
                elif shape.has_table:
                    for row in shape.table.rows:
                        row_cells = " | ".join(cell.text for cell in row.cells)
                        if row_cells.strip():
                            parts.append(row_cells)
        text = "\n".join(parts)
        result = _bounded_text(text, max_chars=max_chars)
        result.backend = "pptx"
        return result
    except Exception:
        return TextExtraction(
            supported=False,
            recoverable=True,
            warnings=("pptx_parse_error",),
        )


def _bounded_text(text: str, *, max_chars: int) -> TextExtraction:
    if len(text) <= max_chars:
        return TextExtraction(supported=True, text=text, truncated=False)
    return TextExtraction(supported=True, text=text[:max_chars], truncated=True)


def _is_pdf(path: Path, mime_type: str) -> bool:
    return mime_type == "application/pdf" or path.suffix.lower() == ".pdf"


def _is_html(path: Path, mime_type: str) -> bool:
    return mime_type == "text/html" or path.suffix.lower() in {".html", ".htm"}


def _is_text_like(path: Path, mime_type: str) -> bool:
    if mime_type.startswith("text/"):
        return True
    return path.suffix.lower() in {".txt", ".md", ".markdown", ".csv"}


def _is_docx(path: Path, mime_type: str) -> bool:
    return (
        mime_type
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        or path.suffix.lower() == ".docx"
    )


def _is_xlsx(path: Path, mime_type: str) -> bool:
    return (
        mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        or path.suffix.lower() == ".xlsx"
    )


def _is_pptx(path: Path, mime_type: str) -> bool:
    return (
        mime_type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or path.suffix.lower() == ".pptx"
    )


def _ensure_size_allowed(size_bytes: int, max_file_size_bytes: int) -> None:
    if size_bytes > max_file_size_bytes:
        raise SlackFileReadError(
            "Slack file is too large to read: "
            f"{size_bytes} bytes exceeds limit {max_file_size_bytes} bytes"
        )


def _content_type(value: object) -> str | None:
    if not isinstance(value, str) or not value.strip():
        return None
    return value.split(";", 1)[0].strip() or None


def _optional_string(value: object) -> str | None:
    if isinstance(value, str) and value.strip():
        return value.strip()
    return None


def _optional_int(value: object) -> int | None:
    if isinstance(value, bool):
        return None
    if isinstance(value, int):
        return value
    if isinstance(value, str) and value.strip().isdigit():
        return int(value)
    return None


class _HTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_starttag(
        self,
        tag: str,
        attrs: list[tuple[str, str | None]],
    ) -> None:
        if tag.lower() in {"br", "div", "li", "p", "tr"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if data:
            self.parts.append(data)

    def text(self) -> str:
        return unescape("".join(self.parts)).strip()
